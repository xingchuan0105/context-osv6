#!/usr/bin/env python3
"""Office 直读解析器（avrag-rs 解析管线 OfficeDirect 后端）。

docx/xlsx/pptx 用直读专用库（mammoth / openpyxl / python-pptx）；旧二进制
doc/ppt/xls 先经 soffice **无损转 OOXML** 再直读（不做 PDF 渲染，避免丢列丢行）。

输出 markdown，供 `blocks_from_markdown` 切成 Heading/Paragraph IR（与 markitdown
路径一致）。

失败语义（hard-fail）：任何失败（含 soffice 转换失败/超时/产物缺失）以非零退出 +
stderr 错误信息返回；**不降级回 markitdown**（对 doc/ppt/xls 是 plain-text 乱码，
等于无兜底）。重试策略由 worker 侧控制。

CLI:
    office-direct-extract <input> <output.md>

环境变量:
    OFFICE_SOFFICE_BIN            默认 `soffice`
    OFFICE_SOFFICE_TIMEOUT_MS     默认 90000（soffice 子进程超时，kill 回收）
"""
from __future__ import annotations

import os
import re
import subprocess
import sys
import tempfile

__all__ = ["convert", "main", "ConverterError"]

#: 纯 hex 文本框（粘贴对象残留，如 `e7d195523061f1c0…`）——pptx 治本跳过。
_HEX_BLOB = re.compile(r"^[0-9A-Fa-f]{40,}$")
#: pptx 限定防御层：长 hex run（git sha 拼接/base64 片段不进代码/文本路径，仅 pptx）。
_HEX_RUN = re.compile(r"[0-9A-Fa-f]{100,}")
#: 空图片占位（mammoth 关内嵌后 markdownify 产 `![]()`）——不进 Paragraph 文本。
_EMPTY_IMG = re.compile(r"!\[\]\(\s*\)")


class ConverterError(Exception):
    """解析失败；message 写入 stderr 供 worker 透传为 ParseWarning。"""


def _soffice_bin() -> str:
    return os.environ.get("OFFICE_SOFFICE_BIN", "soffice")


def _soffice_timeout_s() -> float:
    return int(os.environ.get("OFFICE_SOFFICE_TIMEOUT_MS", "90000")) / 1000.0


def _run_soffice_convert(src: str, out_dir: str, target_ext: str) -> str:
    """soffice --headless --invisible --convert-to <ext> --outdir <dir> <src>。

    - 每次独立临时 UserInstallation profile（LibreOffice 并发 profile 锁竞争）。
    - 失败/超时/产物缺失抛 ConverterError（hard-fail，不降级）。
    - 返回转换产物路径（LibreOffice 会清洗文件名，故扫描目录取唯一匹配）。
    """
    bin_ = _soffice_bin()
    with tempfile.TemporaryDirectory(prefix="office-direct-lo-profile-") as prof:
        cmd = [
            bin_,
            f"-env:UserInstallation=file://{prof}",
            "--headless",
            "--invisible",
            "--convert-to",
            target_ext,
            "--outdir",
            out_dir,
            src,
        ]
        try:
            proc = subprocess.run(
                cmd, capture_output=True, text=True, timeout=_soffice_timeout_s()
            )
        except subprocess.TimeoutExpired as exc:
            raise ConverterError(
                f"soffice timed out after {_soffice_timeout_s()}s: {src}"
            ) from exc
        except FileNotFoundError as exc:
            raise ConverterError(
                f"soffice binary not found: {bin_} "
                "(worker 需装 LibreOffice writer/calc/impress 三组件)"
            ) from exc
        if proc.returncode != 0:
            raise ConverterError(
                f"soffice exited {proc.returncode}: {proc.stderr[-500:]}"
            )
    cand = os.path.join(out_dir, os.path.splitext(os.path.basename(src))[0] + "." + target_ext)
    if os.path.exists(cand):
        return cand
    matches = [f for f in os.listdir(out_dir) if f.endswith("." + target_ext)]
    if not matches:
        raise ConverterError(f"soffice conversion produced no output: {src}")
    return os.path.join(out_dir, matches[0])


def _extract_docx(src: str) -> str:
    import mammoth
    import markdownify

    drop_img = mammoth.images.img_element(lambda image: {"src": ""})
    with open(src, "rb") as f:
        html = mammoth.convert_to_html(f, convert_image=drop_img).value
    md = markdownify.markdownify(html, heading_style="ATX", bullets="-")
    # 空图片占位不进 Paragraph（`![]()`），直读路径不产 embedded_images_json，
    # normalize_parsed_document 的 ImageWithContext 分支不会触发。
    md = _EMPTY_IMG.sub("", md)
    return md


def _extract_xlsx(src: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(src, data_only=True, read_only=False)
    parts: list[str] = []
    for ws in wb.worksheets:
        rows = [
            r
            for r in ws.iter_rows(values_only=True)
            if any(c is not None and str(c).strip() for c in r)
        ]
        if not rows:
            continue
        ncols = max(len(r) for r in rows)
        parts.append(f"## {ws.title}")
        hdr = [str(rows[0][c] or "") for c in range(ncols)]
        parts.append("| " + " | ".join(hdr) + " |")
        parts.append("|" + "---|" * ncols)
        for r in rows[1:]:
            cells: list[str] = []
            for c in range(ncols):
                v = r[c] if c < len(r) else None
                s = "" if v is None else str(v).replace("|", "\\|").replace("\n", "<br>")
                cells.append(s)
            parts.append("| " + " | ".join(cells) + " |")
    return "\n".join(parts)


def _extract_pptx(src: str) -> str:
    from pptx import Presentation

    prs = Presentation(src)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"<!-- Slide {i} -->")
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if not txt:
                    continue
                if _HEX_BLOB.match(txt):
                    continue  # 粘贴残留纯 hex 框（源文档缺陷，跳过）
                parts.append(txt)
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                parts.append(
                    "| " + " | ".join(c.text.strip().replace("\n", " ") for c in tbl.rows[0].cells) + " |"
                )
                parts.append("|" + "---|" * len(tbl.columns))
                for r in range(1, len(tbl.rows)):
                    parts.append(
                        "| "
                        + " | ".join(
                            tbl.rows[r].cells[c].text.strip().replace("\n", " ")
                            for c in range(len(tbl.columns))
                        )
                        + " |"
                    )
        parts.append("")
    md = "\n".join(parts)
    # pptx 限定防御层：仅对 pptx 输出敲长 hex run，绝不套用文本/代码/markitdown 路径。
    md = _HEX_RUN.sub("", md)
    return md


_OOXML_EXT = {"doc": "docx", "ppt": "pptx", "xls": "xlsx"}
_DISPATCH = {"docx": _extract_docx, "xlsx": _extract_xlsx, "pptx": _extract_pptx}


def convert(src: str, out: str) -> int:
    """src → markdown 写入 out，返回字节数。失败抛 ConverterError（hard-fail）。"""
    ext = os.path.splitext(src)[1].lower().lstrip(".")
    if ext in _OOXML_EXT:
        with tempfile.TemporaryDirectory(prefix="office-direct-conv-") as tmp:
            ooxml = _run_soffice_convert(src, tmp, _OOXML_EXT[ext])
            md = _DISPATCH[_OOXML_EXT[ext]](ooxml)
    elif ext in _DISPATCH:
        md = _DISPATCH[ext](src)
    else:
        raise ConverterError(f"unsupported extension: .{ext}")
    with open(out, "w") as f:
        f.write(md)
    return len(md)


def main(argv=None) -> int:
    argv = list(argv if argv is not None else sys.argv[1:])
    if len(argv) != 2:
        print("usage: office-direct-extract <input> <output.md>", file=sys.stderr)
        return 2
    src, out = argv
    try:
        n = convert(src, out)
    except ConverterError as exc:
        print(f"office-direct: {exc}", file=sys.stderr)
        return 1
    print(f"[office-direct] {src} -> {out} ({n} bytes)", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
