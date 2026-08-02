"""office-direct 单测：三格式直读 + hex 跳过 + 空占位 strip + soffice 失败路径。

二进制 doc/ppt/xls → OOXML 回环依赖真实文件（仓库无 fixture），已在 spike 阶段
手动验证（/tmp/pdf_spike/binary_test/）；此处仅测失败路径与直读路径。
"""
import base64
import os
import subprocess
import sys

import pytest

from office_direct.main import ConverterError, convert, main

# 1x1 透明 PNG，用于 docx 图片 fixture。
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGA"
    "hKmMIQAAAABJRU5ErkJggg=="
)


@pytest.fixture
def out_file(tmp_path):
    p = tmp_path / "out.md"
    yield str(p)
    if p.exists():
        p.unlink()


def test_xlsx_multisheet_merged_empty(tmp_path, out_file):
    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "SheetA"
    ws1.append(["编号", "名称", "值"])
    ws1.append([1, "甲", 10])
    ws1.append([])  # 全空行 → 应被跳过
    ws1.merge_cells("A4:B4")
    ws1["A4"] = "合并"
    ws2 = wb.create_sheet("SheetB")
    ws2.append(["x"])
    src = tmp_path / "t.xlsx"
    wb.save(src)

    convert(str(src), out_file)
    md = open(out_file).read()
    assert "## SheetA" in md and "## SheetB" in md
    assert "| 编号 | 名称 | 值 |" in md
    assert "| 1 | 甲 | 10 |" in md
    assert "| 合并 |  |  |" in md  # 合并取左上值
    assert "|  |  |  |" not in md  # 空行不进管道表


def test_docx_heading_table_paragraph(tmp_path, out_file):
    import docx

    d = docx.Document()
    d.add_heading("一 绪论", level=1)
    d.add_heading("1.1 背景", level=2)
    d.add_paragraph("这是正文段落，包含关键短语。")
    tbl = d.add_table(rows=2, cols=3)
    tbl.rows[0].cells[0].text = "措施"
    tbl.rows[0].cells[1].text = "内容"
    tbl.rows[1].cells[0].text = "健全冷链物流标准和服务规范体系"
    tbl.rows[1].cells[1].text = "系统梳理修订完善"
    src = tmp_path / "t.docx"
    d.save(str(src))

    convert(str(src), out_file)
    md = open(out_file).read()
    assert "# 一 绪论" in md
    assert "## 1.1 背景" in md
    assert "这是正文段落" in md
    assert "健全冷链物流标准和服务规范体系" in md  # 单元格完整，不被拆碎


def test_docx_empty_image_placeholder_stripped(tmp_path, out_file):
    import docx

    d = docx.Document()
    d.add_paragraph("图前文字")
    img = tmp_path / "px.png"
    img.write_bytes(_PNG)
    d.add_picture(str(img))
    d.add_paragraph("图后文字")
    src = tmp_path / "t.docx"
    d.save(str(src))

    convert(str(src), out_file)
    md = open(out_file).read()
    assert "![]()" not in md  # 空图片占位被敲掉
    assert "![]( )" not in md
    assert "图前文字" in md and "图后文字" in md


def test_pptx_skips_pure_hex_shape(tmp_path, out_file):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(0, 0, 914400, 914400)
    tb.text_frame.text = "正常文本"
    hb = slide.shapes.add_textbox(0, 914400, 914400, 914400)
    hb.text_frame.text = "e7d195523061f1c0" + "0" * 100  # 纯 hex 残渣
    src = tmp_path / "t.pptx"
    prs.save(str(src))

    convert(str(src), out_file)
    md = open(out_file).read()
    assert "<!-- Slide 1 -->" in md
    assert "正常文本" in md
    assert "e7d195523061f1c0" not in md  # 纯 hex 形状被跳过


def test_pptx_hex_run_defensive_strip(tmp_path, out_file):
    """pptx 限定防御层：长 hex run 被敲，短 hex（合法）保留。"""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(0, 0, 914400, 914400)
    tb.text_frame.text = "债券编号 000002 与 " + "A" * 120 + " 结尾"
    src = tmp_path / "t.pptx"
    prs.save(str(src))

    convert(str(src), out_file)
    md = open(out_file).read()
    assert "000002" in md  # 短 hex 合法内容保留
    assert "A" * 120 not in md  # 长 hex run 被敲


def test_soffice_missing_binary_hard_fail(tmp_path, out_file):
    src = tmp_path / "t.doc"
    src.write_bytes(b"OLE binary stub")  # 内容无关，未到转换即失败
    os.environ["OFFICE_SOFFICE_BIN"] = "office-direct-no-such-soffice"
    try:
        with pytest.raises(ConverterError, match="soffice binary not found"):
            convert(str(src), out_file)
    finally:
        os.environ.pop("OFFICE_SOFFICE_BIN", None)
    assert not os.path.exists(out_file) or os.path.getsize(out_file) == 0


def test_unsupported_extension(tmp_path, out_file):
    src = tmp_path / "t.xyz"
    src.write_bytes(b"x")
    with pytest.raises(ConverterError, match="unsupported extension"):
        convert(str(src), out_file)


def test_main_exit_codes(tmp_path):
    src = tmp_path / "t.docx"
    src.write_bytes(b"not a real docx")
    assert main([]) == 2
    assert main(["a", "b", "c"]) == 2
    # 非 docx/xlsx/pptx 扩展名 → hard fail 退出 1
    assert main([str(tmp_path / "t.xyz"), str(tmp_path / "o.md")]) == 1
